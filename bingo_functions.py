import time
import pyautogui
import pyttsx3
import speech_recognition as sr
import os
from datetime import datetime
import spotipy
from spotipy.oauth2 import SpotifyOAuth
import webbrowser
import pywhatkit as kit
import pandas as pd
import requests
import random
from dotenv import load_dotenv
from instabot import Bot
import instaloader
import pdfplumber
import docx
import matplotlib.pyplot as plt
from pptx import Presentation
from collections import Counter
from nltk.tokenize import sent_tokenize
from nltk.corpus import stopwords
from nltk.probability import FreqDist
import nltk


nltk.download('punk')
nltk.download('stopwords')
nltk.download('punkt_tab')

load_dotenv()

SPOTIPY_CLIENT_ID = os.getenv("SPOTIPY_CLIENT_ID")
SPOTIPY_CLIENT_SECRET = os.getenv("SPOTIPY_CLIENT_SECRET")
SPOTIPY_REDIRECT_URI = os.getenv("SPOTIPY_REDIRECT_URI")
NEWS_API_KEY = os.getenv("NEWS_API_KEY")
SENDER_EMAIL = os.getenv("SENDER_EMAIL")
SENDER_PASSWORD = os.getenv("SENDER_PASSWORD")
IG_USERNAME = os.getenv("IG_USERNAME")
IG_PASSWORD = os.getenv("IG_PASSWORD")


file_path = './contact.csv'
contacts = pd.read_csv(file_path, sep=';', on_bad_lines='skip', engine='python')
contacts_filtered = contacts[['Name', 'Contact']]
contacts_filtered = contacts_filtered.dropna(subset=['Name', 'Contact'])
contact_dict = pd.Series(contacts_filtered['Contact'].values, index=contacts_filtered['Name'].str.lower()).to_dict()

engine = pyttsx3.init('nsss')
voices = engine.getProperty('voices')
print(voices[1].id)
engine.setProperty('voices', voices[0].id)

#text to audio
def speak(audio):
    engine.say(audio)
    print(audio)
    engine.runAndWait()

# voice to text
def takecommand():
    try:
        recognizer = sr.Recognizer()
        os.system("osascript -e 'set volume input volume 100'")
        with sr.Microphone(device_index=0) as source:
            print("Adjusting for ambient noise... Please wait.")
            recognizer.adjust_for_ambient_noise(source, duration=1)
            print("Listening...")
            audio = recognizer.listen(source, phrase_time_limit=5)

        try:
            print("Recognizing...")
            query = recognizer.recognize_google(audio, language='en-in')
            print(f"User said: {query}")
            return query

        except sr.UnknownValueError:
            print("Could not understand the audio. Please speak clearly.")
            return "none"
        except sr.RequestError as e:
            print(f"Recognition request failed; {e}")
            return "none"

    except Exception as e:
        print(f"An error occurred: {e}")
        return "none"


def wish():
    now = datetime.now()
    hour = now.hour
    minute = now.minute

    if 0 <= hour < 12:
        speak("Good morning!")
        speak(f"The time is {hour}:{minute:02d} AM.")
    elif 12 <= hour < 18:
        speak("Good afternoon!")
        speak(f"The time is {hour - 12 if hour > 12 else hour}:{minute:02d} PM.")
    else:
        speak("Good evening!")
        speak(f"The time is {hour - 12 if hour > 12 else hour}:{minute:02d} PM.")

    speak("I am Bingo, here to assist you. How can I help you today?")

def open_spotify_web_with_track(track_uri):
    track_id = track_uri.split(":")[-1]
    web_url = f"https://open.spotify.com/track/{track_id}"
    try:
        webbrowser.open(web_url)
    except Exception as e:
        print(f"Could not open Spotify Web Player: {e}")

def play_random_song():
    # Search for a popular playlist and play a random track from it
    results = sp.search(q='Top Hits', type='playlist', limit=1)
    if results['playlists']['items']:
        playlist_uri = results['playlists']['items'][0]['uri']
        playlist_id = playlist_uri.split(":")[-1]
        tracks = sp.playlist_tracks(playlist_id)['items']

        if tracks:
            random_track = random.choice(tracks)
            track_uri = random_track['track']['uri']
            print(f"Playing {random_track['track']['name']} by {random_track['track']['artists'][0]['name']}...")
            open_spotify_web_with_track(track_uri)
            print("Please ensure you are logged into Spotify Web Player and press play.")
        else:
            print("No tracks found in the playlist.")
    else:
        print("Playlist not found. Please try again.")


SPOTIPY_REDIRECT_URI = 'http://localhost:8898/callback'
scope = 'user-read-playback-state'

sp = spotipy.Spotify(auth_manager=SpotifyOAuth(
    client_id=SPOTIPY_CLIENT_ID,
    client_secret=SPOTIPY_CLIENT_SECRET,
    redirect_uri=SPOTIPY_REDIRECT_URI,
    scope=scope
))

def open_spotify_web_with_track(track_uri):
    track_id = track_uri.split(":")[-1]
    web_url = f"https://open.spotify.com/track/{track_id}"
    try:
        webbrowser.open(web_url)
    except Exception as e:
        print(f"Could not open Spotify Web Player: {e}")

def play_random_song():
    # Search for a popular playlist and play a random track from it
    results = sp.search(q='Top Hits', type='playlist', limit=1)
    if results['playlists']['items']:
        playlist_uri = results['playlists']['items'][0]['uri']
        playlist_id = playlist_uri.split(":")[-1]
        tracks = sp.playlist_tracks(playlist_id)['items']

        if tracks:
            random_track = random.choice(tracks)
            track_uri = random_track['track']['uri']
            print(f"Playing {random_track['track']['name']} by {random_track['track']['artists'][0]['name']}...")
            open_spotify_web_with_track(track_uri)
            print("Please ensure you are logged into Spotify Web Player and press play.")
        else:
            print("No tracks found in the playlist.")
    else:
        print("Playlist not found. Please try again.")
def format_phone_number(phone_number):
    if not phone_number.startswith('+'):
        phone_number = '+' + phone_number
    return phone_number


def play_random_song_on_youtube(search_query="Top hits songs 2024"):
    video_url = kit.playonyt(search_query, open_video=False)
    if video_url:
        autoplay_url = video_url + "&autoplay=1"
        webbrowser.open(autoplay_url)
    else:
        print("Could not find a video to play.")

def fetch_news():
    url = f"https://newsapi.org/v2/everything?q=India&language=en&apiKey=={NEWS_API_KEY}"
    response = requests.get(url)
    news_data = response.json()

    if news_data["status"] == "ok" and news_data["totalResults"] > 0:
        articles = news_data["articles"]  # Get all articles
        headlines = [article["title"] for article in articles]
        return headlines
    else:
        return ["No news articles found at the moment."]

_instabot_instance = None

def get_instabot_instance():
    global _instabot_instance
    if _instabot_instance is None:
        _instabot_instance = Bot()
        _instabot_instance.login(username=os.getenv("IG_USERNAME"), password=os.getenv("IG_PASSWORD"))
        print("Logged in to Instagram successfully!")
    return _instabot_instance


def upload_to_instagram():
    bot = get_instabot_instance()

    try:
        image_paths = input("Enter image file paths separated by commas: ").split(",")
        image_paths = [path.strip() for path in image_paths if path.strip()]
        speak("dictate caption")
        caption = takecommand().lower()
        if len(image_paths) > 1:
            bot.upload_album(image_paths, caption=caption)
        else:
            bot.upload_photo(image_paths[0], caption=caption)

        print("Image(s) uploaded to Instagram.")
    except Exception as e:
        print(f"Failed to upload image(s) to Instagram: {e}")


def location():
    speak("Wait boss, let me check.")
    try:
        # Get the IP address
        ipAdd = requests.get('https://api.ipify.org').text
        print("IP Address:", ipAdd)

        # Get the geographical data
        url = f'https://get.geojs.io/v1/ip/geo/{ipAdd}.json'
        geo_requests = requests.get(url)
        geo_data = geo_requests.json()

        # Extract location details
        city = geo_data.get('city', 'unknown city')
        state = geo_data.get('region', 'unknown state')
        country = geo_data.get('country', 'unknown country')

        # Announce the location
        speak(f"Sir, I think we are near  {city} city, in {state} state of {country}.")
    except Exception as e:
        speak("Sorry sir, due to a network issue I am not sure.")
        print(f"Error: {e}")


def check_insta_id():
    speak("Sir, please enter the username.")
    name = input("Enter the username: ")

    # Open the Instagram profile in the browser
    webbrowser.open(f"https://www.instagram.com/{name}")
    speak(f"Sir, here is the profile of {name}.")

    time.sleep(5)  # Wait before asking to download

    speak("Would you like to download the profile picture?")
    condition = takecommand().lower()

    if "yes" in condition:
        mod = instaloader.Instaloader()
        mod.download_profile(name, profile_pic_only=True)
        speak("Downloaded profile picture.")
    else:
        speak("Alright, not downloading the profile picture.")


def take_screenshot():
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    filename = f"screenshot_{timestamp}.png"
    screenshot = pyautogui.screenshot()
    screenshot.save(filename)
    speak(f"Screenshot taken and saved as {filename}")
    print(f"Screenshot saved as {filename}")


# Summarization and Text Analysis
def summarize_text(text: str) -> dict:
    sentences = sent_tokenize(text)
    word_count = len(text.split())
    stop_words = set(stopwords.words("english"))
    words = [word for word in text.split() if word.lower() not in stop_words]
    freq_dist = FreqDist(words)
    most_common_words = freq_dist.most_common(10)
    summary = " ".join(sentences[:3]) if len(sentences) > 3 else text
    return {
        'summary': summary,
        'total_words': word_count,
        'most_common_words': most_common_words
    }

def analyze_csv(file_path: str) -> dict:
    df = pd.read_csv(file_path)
    numeric_columns = df.select_dtypes(include='number').columns
    text = "CSV Analysis: Contains {} rows and {} columns. Numerical columns are {}.".format(df.shape[0], df.shape[1], ", ".join(numeric_columns))
    for col in numeric_columns:
        plt.figure()
        df[col].hist()
        plt.title(f'Histogram of {col}')
        plt.xlabel(col)
        plt.ylabel('Frequency')
        plt.show()
    return {
        'summary': text,
        'numeric_description': df.describe().to_dict()
    }

def read_txt(file_path: str) -> str:
    with open(file_path, 'r') as file:
        return file.read()

def read_pdf(file_path: str) -> str:
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def read_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return " ".join([para.text for para in doc.paragraphs])

def read_pptx(file_path: str) -> str:
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + " "
    return text


def analyze_file(file_path: str) -> dict:
    extension = os.path.splitext(file_path)[1].lower()
    if extension == '.txt':
        text = read_txt(file_path)
        analysis_results = summarize_text(text)
    elif extension == '.pdf':
        text = read_pdf(file_path)
        analysis_results = summarize_text(text)
    elif extension == '.docx':
        text = read_docx(file_path)
        analysis_results = summarize_text(text)
    elif extension == '.csv':
        analysis_results = analyze_csv(file_path)
    elif extension == '.pptx':
        text = read_pptx(file_path)
        analysis_results = summarize_text(text)
    else:
        raise ValueError(f"Unsupported file extension: {extension}")


    return analysis_results

